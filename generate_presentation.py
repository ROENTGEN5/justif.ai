import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

def add_title_slide(title_text, subtitle_text):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    return slide

def add_content_slide(title_text, bullets):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title_text
    
    tf = body_shape.text_frame
    tf.text = bullets[0]
    
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        if bullet.startswith("  -") or bullet.startswith("  *"):
            p.level = 1
            p.text = bullet.strip(" -*")

# Slide 1: Title
add_title_slide(
    "Securing Justif.ai: Role-Based Access Control (RBAC)",
    "Backend Authentication & Authorization Architecture\n\nProject: Justif.ai (\"Ang iyong munting Abogado\")\nFocus: FastAPI Middleware & Supabase JWT Integration"
)

# Slide 2
add_content_slide(
    "How Authentication Works in Justif.ai",
    [
        "Frontend (Client): Mobile app authenticates with Supabase Auth.",
        "Token Handshake: Supabase returns a JSON Web Token (JWT) containing user identity and role.",
        "Backend (FastAPI): Every protected API request includes the JWT in the Authorization: Bearer header.",
        "Middleware (auth.py): Intercepts the request, cryptographically verifies the token, and extracts the user's role before allowing database access."
    ]
)

# Slide 3
add_content_slide(
    "The AuthenticatedUser Model",
    [
        "Once a token is verified, the system creates an AuthenticatedUser object.",
        "Key Attributes Extracted:",
        "  - user_id (UUID): Links the request to specific database records.",
        "  - email: For communication and auditing.",
        "  - role: Determines permissions (defaults to 'authenticated').",
        "  - token: The raw JWT, passed down to Supabase for secure Row-Level Security (RLS) execution."
    ]
)

# Slide 4
add_content_slide(
    "Dual-Algorithm Token Verification",
    [
        "Our system is built to handle modern infrastructure transitions seamlessly.",
        "ES256 (Asymmetric - Modern Supabase):",
        "  - Verifies tokens using Supabase's public keys via a JWKS (JSON Web Key Set) endpoint.",
        "HS256 (Symmetric - Legacy/Fallback):",
        "  - Verifies tokens using the private SUPABASE_JWT_SECRET stored securely in backend environment variables.",
        "Benefit: Zero downtime regardless of how the Supabase project is configured."
    ]
)

# Slide 5
add_content_slide(
    "Extracting the RBAC Claims",
    [
        "The JWT payload contains 'claims' (trusted data from the auth server).",
        "The Process:",
        "  - 1. Token is cryptographically decoded.",
        "  - 2. The 'sub' claim becomes the user_id.",
        "  - 3. The 'role' claim dictates system privileges.",
        "Failsafe Mechanism: If no role is present, the system defaults to the least privileged standard role: 'authenticated'.",
        "Security Check: If the 'sub' claim is missing, the request is instantly rejected with a 401 Unauthorized."
    ]
)

# Slide 6
add_content_slide(
    "Protecting FastAPI Routes",
    [
        "Routes are secured using FastAPI's Dependency Injection (Depends(get_current_user)).",
        "Example Usage:",
        "  - @router.post('/chat')",
        "  - async def send_message(user: AuthenticatedUser = Depends(get_current_user)): ...",
        "Extensibility: We can easily create dependencies like require_admin_role that check if user.role == 'admin' before allowing access to sensitive administrative endpoints."
    ]
)

# Slide 7
add_content_slide(
    "Defense in Depth: Database Authorization",
    [
        "The backend does not rely purely on API-level RBAC.",
        "By passing the raw token inside the AuthenticatedUser model, we initialize a localized Supabase client per request.",
        "Result: Supabase's PostgreSQL database enforces Row Level Security (RLS).",
        "Even if an API route is compromised, the database still checks the user's role and ID before allowing Data Manipulation (Insert/Select/Delete)."
    ]
)

# Slide 8
add_content_slide(
    "Summary & Next Steps",
    [
        "Key Takeaways:",
        "  - Secure, dual-algorithm token verification.",
        "  - Extensible AuthenticatedUser model.",
        "  - Defense-in-depth architecture integrating FastAPI dependency injection with PostgreSQL RLS.",
        "Next Steps:",
        "  - Define custom roles in Supabase (e.g., premium_user, admin).",
        "  - Add custom FastAPI dependencies for role-specific route protection."
    ]
)

output_file = r"C:\Users\Mark\Documents\courses\Capstone_Project_1\Justif.ai\Justifai_RBAC_Presentation.pptx"
prs.save(output_file)
print(f"Presentation saved successfully to {output_file}")
